#!/usr/bin/env python3
"""
Build defense presentation for LVTN:
"Thiết kế RTL và xây dựng môi trường kiểm thử UVM cho lõi vi xử lý RISC-V (RV32IM) có hỗ trợ ngắt"
SVTH: Phạm Tiến Khoa — GVHD: Trần Hoàng Linh
"""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

TEMPLATE = "/tmp/pp-maulvtn_2026.pptx"
OUTPUT = "/home/key/workspace/project2_axi/doc/Thesis/defense_presentation.pptx"

# Colors from template
BLUE_TITLE = RGBColor(0x33, 0x76, 0xC7)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE = RGBColor(0x1F, 0x38, 0x96)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x10, 0x89, 0x42)
RED = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation(TEMPLATE)
W = prs.slide_width   # 9144000 EMU
H = prs.slide_height  # 6858000 EMU

# ── helpers ────────────────────────────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                word_wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items,
                       font_size=15, color=BLACK, bullet_char="▪ ",
                       indent_items=None, title=None, title_size=17,
                       title_color=BLUE_TITLE):
    """Add a textbox with bullet list. indent_items: list of (level, text)."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    first_para = True
    if title:
        p = tf.paragraphs[0] if first_para else tf.add_paragraph()
        first_para = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color

    for item in items:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item

        p = tf.paragraphs[0] if first_para else tf.add_paragraph()
        first_para = False
        p.alignment = PP_ALIGN.LEFT
        indent = "    " * level
        run = p.add_run()
        run.text = f"{indent}{bullet_char}{text}"
        run.font.size = Pt(font_size - level * 1)
        run.font.bold = False
        run.font.color.rgb = color

    return txBox


def set_slide_title(slide, text, font_size=24, color=BLUE_TITLE):
    """Set the title placeholder text."""
    for shape in slide.shapes:
        if shape.shape_type == 14:  # PLACEHOLDER
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = color
            return shape
    return None


def add_slide_from_layout(layout_idx=2):
    """Add a new slide from a given slide layout index."""
    layout = prs.slide_layouts[layout_idx] if layout_idx < len(prs.slide_layouts) else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    # Remove placeholder hints
    for shape in list(slide.placeholders):
        if shape.has_text_frame:
            shape.text_frame.clear()
    return slide


def clear_content_shapes(slide):
    """Remove non-placeholder, non-slide-number shapes (cleanup template notes)."""
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 17:  # TEXT_BOX
            to_remove.append(shape)
        elif shape.shape_type == 1 and "TextBox 47" in shape.name:
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_divider_line(slide, top, color_rgb=(0x33, 0x76, 0xC7)):
    """Add a thin horizontal line."""
    from pptx.util import Pt as PtU
    line = slide.shapes.add_connector(
        1,  # STRAIGHT
        Emu(685440), Emu(top),
        Emu(685440 + 7772640), Emu(top)
    )
    line.line.color.rgb = RGBColor(*color_rgb)
    line.line.width = Pt(1.5)


# ── Slide 1: Title ──────────────────────────────────────────────────────────

slide1 = prs.slides[0]
clear_content_shapes(slide1)

# School info top
add_textbox(slide1,
    left=0, top=180000, width=W, height=350000,
    text="ĐẠI HỌC QUỐC GIA TP. HỒ CHÍ MINH — TRƯỜNG ĐẠI HỌC BÁCH KHOA",
    font_size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_textbox(slide1,
    left=0, top=480000, width=W, height=280000,
    text="Khoa Điện – Điện Tử  |  Bộ Môn Điện Tử",
    font_size=12, bold=False, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Title box
add_textbox(slide1,
    left=400000, top=900000, width=W - 800000, height=1600000,
    text="THIẾT KẾ RTL VÀ XÂY DỰNG MÔI TRƯỜNG KIỂM THỬ UVM\nCHO LÕI VI XỬ LÝ RISC-V (RV32IM) CÓ HỖ TRỢ NGẮT",
    font_size=24, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide1,
    left=0, top=2650000, width=W, height=280000,
    text="LUẬN VĂN TỐT NGHIỆP ĐẠI HỌC — 2026",
    font_size=14, bold=False, color=BLUE_TITLE, align=PP_ALIGN.CENTER)

# SVTH / GVHD
add_textbox(slide1,
    left=1800000, top=3200000, width=W - 3600000, height=900000,
    text="SVTH: Phạm Tiến Khoa     MSSV: 2010344\nGVHD: TS. Trần Hoàng Linh\nBộ Môn Điện Tử — Khoa Điện – Điện Tử",
    font_size=15, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

add_textbox(slide1,
    left=0, top=4400000, width=W, height=260000,
    text="TP. HỒ CHÍ MINH, THÁNG 5 NĂM 2026",
    font_size=13, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

print("✓ Slide 1 (Title) done")

# ── Slide 2: Nội dung (Table of Contents) ──────────────────────────────────

slide2 = prs.slides[1]
clear_content_shapes(slide2)
set_slide_title(slide2, "NỘI DUNG TRÌNH BÀY")

toc_items = [
    "1.  Giới thiệu tổng quan",
    "2.  Nhiệm vụ luận văn",
    "3.  Cơ sở lý thuyết",
    "4.  Thiết kế phần cứng RTL (RV32IM Pipeline)",
    "5.  Xây dựng môi trường kiểm thử UVM",
    "6.  Kết quả thực hiện",
    "7.  Kết luận và hướng phát triển",
]
add_bullet_textbox(slide2,
    left=685440, top=1050000, width=7772640, height=5400000,
    items=toc_items, font_size=20, bullet_char="", color=BLACK)

print("✓ Slide 2 (TOC) done")

# ── Slide 3: Giới thiệu tổng quan ──────────────────────────────────────────

slide3 = prs.slides[2]
clear_content_shapes(slide3)
set_slide_title(slide3, "1. Giới thiệu tổng quan")

add_bullet_textbox(slide3,
    left=685440, top=1000000, width=7772640, height=2200000,
    title="Bối cảnh",
    items=[
        "RISC-V — kiến trúc tập lệnh mã nguồn mở, đang phổ biến mạnh trong IC Design toàn cầu",
        "Design Verification (DV) chiếm tới 70% tổng thời gian & nguồn lực trong quy trình thiết kế chip",
        "UVM (IEEE 1800.2) là chuẩn công nghiệp hàng đầu cho Functional Verification",
    ], font_size=16, color=BLACK)

add_bullet_textbox(slide3,
    left=685440, top=3400000, width=7772640, height=2800000,
    title="Mục tiêu đề tài",
    items=[
        "Thiết kế lõi vi xử lý RV32IM pipeline 5 tầng hoàn chỉnh, có hỗ trợ ngắt (M-mode)",
        "Áp dụng UVM chuẩn công nghiệp để xây dựng môi trường kiểm thử tự động",
        "Đạt 100% Functional Coverage, kiểm thử >10,000 lệnh ngẫu nhiên",
    ], font_size=16, color=BLACK)

print("✓ Slide 3 (Introduction) done")

# ── Slide 4: Nhiệm vụ luận văn ──────────────────────────────────────────────

slide4 = prs.slides[3]
clear_content_shapes(slide4)
set_slide_title(slide4, "2. Nhiệm vụ luận văn")

add_bullet_textbox(slide4,
    left=685440, top=950000, width=7772640, height=5400000,
    title="Các nhiệm vụ cụ thể:",
    items=[
        "NV 1: Nghiên cứu đặc tả RISC-V (Unprivileged ISA + Privileged Architecture) và chuẩn UVM IEEE 1800.2",
        "NV 2: Thiết kế RTL bằng SystemVerilog cho lõi RV32IM pipeline 5 tầng\n"
            + "         → 47 lệnh RV32IM + 6 lệnh CSR + MRET/ECALL/EBREAK\n"
            + "         → Forwarding, Load-Use Stall, Branch Flush, CSR, Interrupt (M-mode)",
        "NV 3: Xây dựng môi trường kiểm thử UVM\n"
            + "         → Unit Test: ALU, M-Unit, Decoder, LSU\n"
            + "         → Integration Test: Hazard, Interrupt, Stress Test",
        "NV 4: Chạy mô phỏng trên QuestaSim, phân tích kết quả,\n"
            + "         đo lường Functional Coverage và viết báo cáo",
    ], font_size=15, color=BLACK)

print("✓ Slide 4 (Tasks) done")

# ── Slide 5: Cơ sở lý thuyết ────────────────────────────────────────────────

slide5 = prs.slides[4]
clear_content_shapes(slide5)
set_slide_title(slide5, "3. Cơ sở lý thuyết")

add_bullet_textbox(slide5,
    left=685440, top=950000, width=3700000, height=5400000,
    title="Kiến trúc RISC-V",
    items=[
        "32 thanh ghi 32-bit (x0 hardwired = 0)",
        "6 định dạng lệnh: R, I, S, B, U, J",
        "RV32I: 37 lệnh cơ sở",
        "RV32M: 8 lệnh nhân/chia",
        "Divide-by-zero: không raise exception",
        "Privileged: M-mode + CSR registers",
        "Ngắt: Software, Timer, External",
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide5,
    left=4700000, top=950000, width=3757880, height=5400000,
    title="Chuẩn UVM (IEEE 1800.2)",
    items=[
        "uvm_test → uvm_env → uvm_agent",
        "Driver: tạo stimulus (lệnh → IMEM)",
        "Monitor: thu thập tín hiệu DUT",
        "Scoreboard + Reference Model",
        "Sequence / Sequence Item (Transaction)",
        "Constrained-Random Verification",
        "Functional Coverage (Covergroup)",
        "UVM Phases: build → connect → run → report",
    ], font_size=14, color=BLACK)

print("✓ Slide 5 (Theory) done")

# ── Slide 6: Thiết kế phần cứng — Kiến trúc ────────────────────────────────

slide6 = prs.slides[5]
clear_content_shapes(slide6)
set_slide_title(slide6, "4. Thiết kế phần cứng — Kiến trúc tổng quan")

add_bullet_textbox(slide6,
    left=685440, top=950000, width=7772640, height=1300000,
    title="Yêu cầu thiết kế",
    items=[
        "ISA: RV32I + RV32M  |  Pipeline: 5 tầng (IF→ID→EX→MEM→WB)  |  Privilege: M-mode",
        "Ngôn ngữ: SystemVerilog IEEE 1800-2017  |  Giao thức bộ nhớ: Handshake Valid/Ready (AXI-like)",
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide6,
    left=685440, top=2400000, width=7772640, height=3900000,
    title="Cấu trúc module (src/)",
    items=[
        "soc_top.sv  →  riscv_datapath.sv  +  csr.sv  +  riscv_control.sv",
        (1, "pc_gen.sv — PC register + next-PC MUX"),
        (1, "decoder.sv — Giải mã 56 lệnh, tạo control bundle id_ctrl"),
        (1, "alu.sv — 11 phép toán ALU  +  riscv_m_unit.sv — MUL(2 cy)/DIV(35 cy)"),
        (1, "forwarding_unit.sv — MEM→EX priority over WB→EX"),
        (1, "hazard_unit.sv — Load-Use Stall (1 cy)"),
        (1, "lsu.sv — FSM 4 trạng thái, Handshake Valid/Ready, Byte Enable"),
        (1, "regfile.sv — 32×32-bit, x0 hardwired, 2 read / 1 write"),
        (1, "csr.sv — 8 CSR cốt lõi: mstatus, mie, mip, mtvec, mepc, mcause, mtval, mscratch"),
    ], font_size=13, color=BLACK)

print("✓ Slide 6 (HW Architecture) done")

# ── Slide 7: Thiết kế phần cứng — Xử lý Hazard & Ngắt ──────────────────────

# Add new slide by copying slide 6's layout
new_slide7_layout = slide6.slide_layout
slide7 = prs.slides.add_slide(new_slide7_layout)
for shape in list(slide7.placeholders):
    if shape.has_text_frame:
        shape.text_frame.clear()

set_slide_title(slide7, "4. Thiết kế phần cứng — Xử lý Hazard & Ngắt")

add_bullet_textbox(slide7,
    left=685440, top=950000, width=3700000, height=2600000,
    title="Xử lý Data Hazard",
    items=[
        "Data Forwarding (2 đường):",
        (1, "MEM→EX: ưu tiên cao hơn WB→EX"),
        (1, "WB→EX: fallback nếu MEM không match"),
        "Load-Use Stall: 1 chu kỳ NOP",
        (1, "hazard_unit phát hiện is_load_use"),
        (1, "ctrl_force_stall_id = 1"),
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide7,
    left=685440, top=3700000, width=3700000, height=2700000,
    title="Xử lý Control Hazard",
    items=[
        "Branch Resolution tại tầng EX (Late)",
        "Branch Penalty: 2 chu kỳ",
        "Flush IF/ID và ID/EX khi branch_taken=1",
        "JAL: cũng 2 chu kỳ penalty",
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide7,
    left=4700000, top=950000, width=3757880, height=5400000,
    title="Xử lý Ngắt & CSR (M-mode)",
    items=[
        "3 nguồn ngắt: Software, Timer, External",
        "Điều kiện chấp nhận ngắt:",
        (1, "mstatus.MIE = 1 (bật ngắt toàn cục)"),
        (1, "mie[bit] = 1 (bật ngắt cụ thể)"),
        (1, "mip[bit] = 1 (ngắt đang pending)"),
        "Quy trình Trap:",
        (1, "mepc ← PC_current"),
        (1, "mcause ← trap_cause"),
        (1, "mstatus.MIE ← 0 (khóa ngắt ngay)"),
        (1, "PC ← mtvec (ISR address)"),
        "Thoát: MRET khôi phục mstatus.MIE và PC ← mepc",
    ], font_size=13, color=BLACK)

print("✓ Slide 7 (Hazard & Interrupt) done")

# ── Slide 8: UVM Testbench ──────────────────────────────────────────────────

slide8_layout = slide6.slide_layout
slide8 = prs.slides.add_slide(slide8_layout)
for shape in list(slide8.placeholders):
    if shape.has_text_frame:
        shape.text_frame.clear()

set_slide_title(slide8, "5. Xây dựng môi trường kiểm thử UVM")

add_bullet_textbox(slide8,
    left=685440, top=950000, width=3700000, height=5400000,
    title="Các thành phần UVM",
    items=[
        "riscv_transaction (Sequence Item)",
        (1, "opcode, rs1, rs2, rd, imm, result_exp"),
        (1, "Constrained-random với rand + constraint"),
        "riscv_driver",
        (1, "Nạp chương trình vào IMEM ảo"),
        (1, "Cung cấp DMEM Responder với độ trễ ngẫu nhiên 1-5 cy"),
        "riscv_monitor",
        (1, "Sample write-back: rf_we, rd_addr, wdata"),
        "riscv_scoreboard",
        (1, "Reference Model tính toán kết quả kỳ vọng"),
        (1, "So sánh exp_fifo vs act_fifo"),
        "riscv_coverage",
        (1, "Covergroup cho tất cả 56 lệnh + hazard scenarios"),
    ], font_size=13, color=BLACK)

add_bullet_textbox(slide8,
    left=4700000, top=950000, width=3757880, height=5400000,
    title="Kịch bản kiểm thử",
    items=[
        "Unit Test (Đơn vị):",
        (1, "ALU: 11 phép toán, 1000 test vectors"),
        (1, "M-Unit: MUL/DIV, corner cases RISC-V spec"),
        (1, "Decoder: 56 lệnh, 9 lệnh hệ thống"),
        (1, "LSU: LB/LH/LW/LBU/LHU + SB/SH/SW, misaligned"),
        "Integration Test (Toàn hệ thống):",
        (1, "Data Hazard & Forwarding (EX→EX, MEM→EX)"),
        (1, "Load-Use Stall (1 chu kỳ NOP)"),
        (1, "Control Hazard & Pipeline Flush"),
        (1, "Interrupt + CSR (mtvec → ISR → MRET)"),
        "Stress Test:",
        (1, "100 lệnh ngẫu nhiên × 10 lần = 1,000+ lệnh/run"),
        (1, "50% ALU-RAW + 17% Load-Use + 17% Branch + 16% Store"),
    ], font_size=13, color=BLACK)

print("✓ Slide 8 (UVM) done")

# ── Slide 9: Kết quả ────────────────────────────────────────────────────────

slide9_src = prs.slides[7]  # "Kết quả" template slide
clear_content_shapes(slide9_src)
set_slide_title(slide9_src, "6. Kết quả thực hiện")

add_bullet_textbox(slide9_src,
    left=685440, top=950000, width=3700000, height=2600000,
    title="Unit Test — Kết quả",
    items=[
        "ALU: 1000/1000 PASS (11 phép toán, biên số 0/MAX/MIN)",
        "M-Unit: PASS — MUL=2cy, DIV=35cy; corner cases RISC-V spec đúng",
        "Decoder: 56 lệnh PASS; illegal_instr=1 đúng; Immediate 6 định dạng đúng",
        "LSU: PASS — Byte Enable, Sign Extension, Misaligned Exception, Handshake",
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide9_src,
    left=685440, top=3700000, width=3700000, height=2700000,
    title="Integration Test — Kết quả",
    items=[
        "Forwarding EX→EX & MEM→EX: Scoreboard PASS",
        "Load-Use Stall 1 chu kỳ: phát hiện đúng is_load_use",
        "Branch Flush 2 chu kỳ: ctrl_flush_if_id + ctrl_flush_id_ex",
        "Interrupt: mepc/mcause đúng, PC→mtvec, MRET khôi phục",
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide9_src,
    left=4700000, top=950000, width=3757880, height=5400000,
    title="Stress Test & Coverage",
    items=[
        "Stress Test (10 lần × 100 lệnh):",
        (1, "10,000+ lệnh ngẫu nhiên — Scoreboard 100% PASS"),
        (1, "Không phát hiện lỗi chức năng nào"),
        (1, "Phát hiện & fix trong quá trình dev:"),
        (2, "Lỗi ưu tiên Forwarding (MEM vs WB)"),
        (2, "Lỗi biên Restoring Division (MIN_INT ÷ -1)"),
        (2, "Lỗi timing xóa MIE khi Trap"),
        "",
        "Functional Coverage:",
        (1, "100% trên 199 coverage bins"),
        (1, "Bao phủ toàn bộ tập lệnh RV32IM"),
        (1, "Bao phủ Data Hazard, Load-Use, Branch, Interrupt"),
        "",
        "Điểm khác biệt so với đề tài tương tự:",
        (1, "Giao thức Valid/Ready chuẩn AXI-like"),
        (1, "UVM với Constrained-Random (không chỉ direct test)"),
    ], font_size=13, color=BLACK)

print("✓ Slide 9 (Results) done")

# ── Slide 10: Kết luận ──────────────────────────────────────────────────────

slide10_src = prs.slides[8]  # "Kết luận" template slide
clear_content_shapes(slide10_src)
set_slide_title(slide10_src, "7. Kết luận và Hướng phát triển")

add_bullet_textbox(slide10_src,
    left=685440, top=950000, width=7772640, height=2800000,
    title="Kết luận — Đã hoàn thành",
    items=[
        "✓ Lõi RV32IM pipeline 5 tầng: 47 lệnh + 9 lệnh hệ thống, SystemVerilog",
        "✓ Hazard handling: Forwarding (MEM/WB→EX), Load-Use Stall, Branch Flush (2cy)",
        "✓ CSR + Ngắt M-mode: 8 CSR cốt lõi, 3 nguồn ngắt, Trap hoàn chỉnh",
        "✓ UVM Testbench: Driver, Monitor, Scoreboard, Coverage — chuẩn công nghiệp",
        "✓ 100% Functional Coverage (199 bins) | Stress Test 10,000+ lệnh — PASS",
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide10_src,
    left=685440, top=3900000, width=3700000, height=2500000,
    title="Hạn chế",
    items=[
        "Chưa tổng hợp lên FPGA (chưa đo Fmax, tài nguyên)",
        "Chưa có PLIC cho đa nguồn ngắt ưu tiên",
        "Branch Penalty 2cy (chưa có Branch Predictor)",
        "Cache chưa được tích hợp",
    ], font_size=14, color=BLACK)

add_bullet_textbox(slide10_src,
    left=4700000, top=3900000, width=3757880, height=2500000,
    title="Hướng phát triển",
    items=[
        "Tích hợp AXI-4 Lite Master Interface",
        "Triển khai FPGA: DE10-Standard / Arty A7",
        "Branch Predictor (BTB + BHT 2-bit)",
        "I-Cache / D-Cache",
        "Formal Verification (SVA): x0=0, FSM deadlock-free",
        "Thêm S-mode, MMU (Sv32) → chạy FreeRTOS/Linux nhẹ",
    ], font_size=14, color=BLACK)

print("✓ Slide 10 (Conclusion) done")

# ── Slide 11: Thank you ─────────────────────────────────────────────────────

slide11_src = prs.slides[9]  # Last "CẢM ƠN" slide
clear_content_shapes(slide11_src)

add_textbox(slide11_src,
    left=0, top=2000000, width=W, height=600000,
    text="XIN CẢM ƠN SỰ LẮNG NGHE",
    font_size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_textbox(slide11_src,
    left=0, top=2800000, width=W, height=400000,
    text="CỦA QUÝ THẦY CÔ VÀ CÁC BẠN",
    font_size=28, bold=False, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_textbox(slide11_src,
    left=0, top=3600000, width=W, height=300000,
    text="Phạm Tiến Khoa — ptk2002ptkt1@gmail.com",
    font_size=16, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

add_textbox(slide11_src,
    left=0, top=4200000, width=W, height=280000,
    text="GitHub: KhoaPham12092002 | MSSV: 2010344",
    font_size=14, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

print("✓ Slide 11 (Thank you) done")

# ── Save ────────────────────────────────────────────────────────────────────

prs.save(OUTPUT)
print(f"\n✅ Saved: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")
