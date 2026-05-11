#!/usr/bin/env python3
"""Build thesis defense presentation for RISC-V RV32IM + UVM thesis."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
import os

THESIS_DIR = "/home/key/workspace/project2_axi/doc/Thesis"
IMG_DIR    = os.path.join(THESIS_DIR, "img")
OUT_FILE   = os.path.join(THESIS_DIR, "thesis_defense.pptx")

# Colors
BK_BLUE   = RGBColor(0x33, 0x76, 0xC7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x26, 0x26, 0x26)
LIGHT_BG  = RGBColor(0xF4, 0xF8, 0xFF)
GREEN_OK  = RGBColor(0x00, 0x96, 0x3A)
ORANGE    = RGBColor(0xE5, 0x7A, 0x00)
RED_ERR   = RGBColor(0xCC, 0x00, 0x00)

W = Inches(10)   # slide width
H = Inches(7.5)  # slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]  # Blank

# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(1)):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_w
    else:
        line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_bullet_box(slide, items, l, t, w, h,
                   font_size=Pt(18), color=DARK_GRAY,
                   bullet_char="▸ ", indent_chars="   "):
    """items: list of (text, level) where level 0=top, 1=sub"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item, level in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2) if level == 0 else Pt(1)
        run = p.add_run()
        prefix = bullet_char if level == 0 else indent_chars + "• "
        run.text = prefix + item
        run.font.size  = font_size if level == 0 else Pt(font_size.pt - 2)
        run.font.color.rgb = color
        run.font.bold = (level == 0)
    return tb

def slide_header(slide, title, section_num=None):
    """Blue header bar + title text."""
    add_rect(slide, 0, 0, W, Inches(0.85), fill_rgb=BK_BLUE)
    display = f"{section_num}. {title}" if section_num else title
    add_text(slide, display,
             Inches(0.25), Inches(0.1), Inches(9.5), Inches(0.7),
             font_size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def slide_footer(slide, page_num):
    """Page number bottom-right."""
    add_text(slide, str(page_num),
             Inches(9.3), Inches(7.15), Inches(0.6), Inches(0.3),
             font_size=Pt(12), color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.RIGHT)

def add_image(slide, img_name, l, t, w, h=None):
    path = os.path.join(IMG_DIR, img_name)
    if not os.path.exists(path):
        # placeholder box
        box = add_rect(slide, l, t, w, h or Inches(2.5),
                       fill_rgb=RGBColor(0xDD,0xEE,0xFF),
                       line_rgb=BK_BLUE)
        add_text(slide, f"[{img_name}]",
                 l + Inches(0.1), t + Inches(0.1), w - Inches(0.2), (h or Inches(2.5)) - Inches(0.2),
                 font_size=Pt(11), color=BK_BLUE, align=PP_ALIGN.CENTER)
        return box
    if h:
        return slide.shapes.add_picture(path, l, t, w, h)
    return slide.shapes.add_picture(path, l, t, w)

def content_area():
    """Return standard content area bounds."""
    return Inches(0.3), Inches(1.0), Inches(9.4), Inches(6.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
# BK Blue top bar
add_rect(slide, 0, 0, W, Inches(1.6), fill_rgb=BK_BLUE)
# White content area
add_rect(slide, 0, Inches(1.6), W, Inches(5.5), fill_rgb=WHITE)
# Bottom blue strip
add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill_rgb=BK_BLUE)

add_text(slide, "ĐẠI HỌC QUỐC GIA TP.HCM — TRƯỜNG ĐẠI HỌC BÁCH KHOA",
         Inches(0), Inches(0.1), W, Inches(0.5),
         font_size=Pt(13), bold=False, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "KHOA ĐIỆN – ĐIỆN TỬ  |  BỘ MÔN ĐIỆN TỬ",
         Inches(0), Inches(0.55), W, Inches(0.4),
         font_size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "LUẬN VĂN TỐT NGHIỆP ĐẠI HỌC",
         Inches(0.5), Inches(1.9), Inches(9), Inches(0.5),
         font_size=Pt(16), bold=False, color=BK_BLUE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Thiết kế RTL và xây dựng môi trường kiểm thử UVM\n"
         "cho lõi vi xử lý RISC-V (RV32IM) có hỗ trợ ngắt",
         Inches(0.5), Inches(2.4), Inches(9), Inches(1.6),
         font_size=Pt(28), bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_text(slide,
         "GVHD:  TS. Trần Hoàng Linh\n"
         "SVTH:   Phạm Tiến Khoa  —  MSSV: 2010344",
         Inches(2.5), Inches(4.3), Inches(5), Inches(1.0),
         font_size=Pt(16), color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "TP. Hồ Chí Minh, 2026",
         Inches(0), Inches(5.5), W, Inches(0.4),
         font_size=Pt(14), color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.CENTER)

slide_footer(slide, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Table of Contents
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=LIGHT_BG)
slide_header(slide, "Nội dung trình bày")
slide_footer(slide, 2)

toc = [
    ("1. Giới thiệu & Động lực đề tài", 0),
    ("2. Nhiệm vụ luận văn", 0),
    ("3. Cơ sở lý thuyết", 0),
    ("   RISC-V ISA (RV32IM) · Pipeline 5 tầng · Hazards · CSR & Ngắt · UVM", 1),
    ("4. Thiết kế phần cứng RTL", 0),
    ("   Datapath · ALU + M-Unit · Forwarding · Load-Use · Interrupt", 1),
    ("5. Thiết kế môi trường kiểm thử UVM", 0),
    ("   Transaction · Driver · Monitor · Scoreboard · Coverage", 1),
    ("6. Kết quả mô phỏng & Phân tích", 0),
    ("7. Kết luận & Hướng phát triển", 0),
]
add_bullet_box(slide, toc, Inches(0.8), Inches(1.0), Inches(8.5), Inches(6.1),
               font_size=Pt(22))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Motivation
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Giới thiệu & Động lực đề tài", "1")
slide_footer(slide, 3)

# Left column — RISC-V context
add_rect(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(5.8),
         fill_rgb=LIGHT_BG, line_rgb=BK_BLUE, line_w=Pt(0.5))
add_text(slide, "Bối cảnh RISC-V", Inches(0.4), Inches(1.05), Inches(4.3), Inches(0.4),
         font_size=Pt(16), bold=True, color=BK_BLUE)
items_l = [
    ("ISA mã nguồn mở, miễn phí bản quyền", 0),
    ("Thiết kế module hóa: RV32I + M + C + F ...", 0),
    ("Được dùng bởi: Google, SiFive, lowRISC", 0),
    ("Xu hướng tất yếu trong IC Design Việt Nam", 0),
    ("", 1),
    ("Quá trình DV chiếm 60–70% tổng nguồn lực", 0),
    ("UVM (IEEE 1800.2) — tiêu chuẩn công nghiệp", 0),
    ("Sinh viên VN hiếm khi tiếp cận UVM", 0),
]
add_bullet_box(slide, items_l, Inches(0.4), Inches(1.5), Inches(4.3), Inches(5.0),
               font_size=Pt(15))

# Right column — Gap in current research
add_rect(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(5.8),
         fill_rgb=RGBColor(0xFF,0xF3,0xE0), line_rgb=ORANGE, line_w=Pt(0.5))
add_text(slide, "Khoảng trống nghiên cứu", Inches(5.3), Inches(1.05), Inches(4.3), Inches(0.4),
         font_size=Pt(16), bold=True, color=ORANGE)
items_r = [
    ("Các đề tài RISC-V tại VN thường dừng ở Directed Testbench viết tay", 0),
    ("Chưa có đề tài kết hợp RTL + UVM chuẩn công nghiệp", 0),
    ("Chưa tích hợp interrupt handling đầy đủ M-mode", 0),
    ("", 1),
    ("⟹ Đề tài này lấp đầy khoảng trống đó", 0),
]
add_bullet_box(slide, items_r, Inches(5.3), Inches(1.5), Inches(4.3), Inches(5.0),
               font_size=Pt(15), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Thesis Tasks
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Nhiệm vụ luận văn", "2")
slide_footer(slide, 4)

tasks = [
    ("Nhiệm vụ 1: Nghiên cứu lý thuyết", 0),
    ("RISC-V Unprivileged ISA (RV32I + M)  ·  RISC-V Privileged Architecture  ·  UVM IEEE 1800.2", 1),
    ("Nhiệm vụ 2: Thiết kế RTL phần cứng (SystemVerilog)", 0),
    ("Pipeline 5 tầng IF/ID/EX/MEM/WB  ·  47 lệnh RV32IM + 9 lệnh hệ thống", 1),
    ("Forwarding Unit, Hazard Detection Unit, CSR Unit, LSU", 1),
    ("Nhiệm vụ 3: Xây dựng môi trường kiểm thử UVM", 0),
    ("Sequence · Driver (IMEM/DMEM Responder) · Monitor · Scoreboard + RefModel", 1),
    ("Constrained Random Verification với 10,000+ lệnh ngẫu nhiên", 1),
    ("Nhiệm vụ 4: Đánh giá kết quả", 0),
    ("Functional Coverage ≥ 95%  ·  Scoreboard 100% PASS  ·  Debug & sửa lỗi RTL", 1),
]
add_bullet_box(slide, tasks, Inches(0.5), Inches(1.0), Inches(9.2), Inches(6.2),
               font_size=Pt(19))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RISC-V ISA Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Cơ sở lý thuyết: Kiến trúc RISC-V RV32IM", "3")
slide_footer(slide, 5)

# Left: ISA description
add_text(slide, "Tập lệnh RV32IM", Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.4),
         font_size=Pt(18), bold=True, color=BK_BLUE)
isa_items = [
    ("RV32I — 47 lệnh số nguyên 32-bit cơ bản", 0),
    ("6 định dạng lệnh: R / I / S / B / U / J", 1),
    ("32 thanh ghi x0–x31  (x0 hardwired = 0)", 1),
    ("RV32M — 8 lệnh nhân/chia phần cứng", 0),
    ("MUL / MULH / MULHSU / MULHU", 1),
    ("DIV / DIVU / REM / REMU", 1),
    ("Chia cho 0: trả về −1 (DIV), không trap", 1),
]
add_bullet_box(slide, isa_items, Inches(0.3), Inches(1.45), Inches(4.5), Inches(4.5),
               font_size=Pt(16))

# Right: 6 formats image
add_image(slide, "6 định dạng RV32I.drawio.png",
          Inches(5.0), Inches(1.0), Inches(4.8), Inches(5.5))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Pipeline 5-stage
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Pipeline 5 tầng & Xung đột đường ống", "3")
slide_footer(slide, 6)

add_text(slide, "Nguyên lý Pipeline: tăng thông lượng (throughput) bằng cách xử lý nhiều lệnh song song theo từng giai đoạn",
         Inches(0.3), Inches(0.95), Inches(9.4), Inches(0.4),
         font_size=Pt(15), color=DARK_GRAY, italic=True)

# Pipeline stages visual
stages = [
    ("IF", "Instruction\nFetch", BK_BLUE),
    ("ID", "Instruction\nDecode", RGBColor(0x00,0x80,0xC0)),
    ("EX", "Execute\nALU/M-Unit", RGBColor(0x00,0xA0,0x60)),
    ("MEM", "Memory\nAccess", ORANGE),
    ("WB", "Write\nBack", RGBColor(0xC0,0x00,0x60)),
]
for idx, (name, label, color) in enumerate(stages):
    x = Inches(0.3 + idx * 1.9)
    add_rect(slide, x, Inches(1.5), Inches(1.7), Inches(1.0), fill_rgb=color)
    add_text(slide, name, x, Inches(1.55), Inches(1.7), Inches(0.4),
             font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, Inches(1.95), Inches(1.7), Inches(0.5),
             font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)
    if idx < 4:
        add_text(slide, "→", Inches(0.3 + idx*1.9 + 1.72), Inches(1.7), Inches(0.18), Inches(0.4),
                 font_size=Pt(20), bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Hazards
add_text(slide, "Ba loại xung đột (Hazards) và giải pháp:",
         Inches(0.3), Inches(2.75), Inches(9.4), Inches(0.35),
         font_size=Pt(17), bold=True, color=BK_BLUE)

hazard_data = [
    ("Xung đột cấu trúc\n(Structural Hazard)", "Harvard architecture:\nIMEM tách DMEM", BK_BLUE),
    ("Xung đột dữ liệu\n(Data Hazard — RAW)", "Data Forwarding (EX/MEM→EX)\n+ 1-cycle Stall (Load-Use)", ORANGE),
    ("Xung đột điều khiển\n(Control Hazard)", "Late Branch Resolution tại EX\n2-cycle Flush (Branch Taken)", RGBColor(0x80,0x00,0x80)),
]
for idx, (title, solution, color) in enumerate(hazard_data):
    x = Inches(0.3 + idx * 3.2)
    add_rect(slide, x, Inches(3.2), Inches(3.0), Inches(1.0), fill_rgb=color)
    add_text(slide, title, x, Inches(3.25), Inches(3.0), Inches(0.85),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, Inches(4.2), Inches(3.0), Inches(2.0),
             fill_rgb=LIGHT_BG, line_rgb=color, line_w=Pt(1))
    add_text(slide, solution, x + Inches(0.1), Inches(4.3), Inches(2.8), Inches(1.8),
             font_size=Pt(14), color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CSR & Interrupt
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Cơ chế xử lý ngắt M-mode (RISC-V Privileged)", "3")
slide_footer(slide, 7)

# CSR Registers table
add_text(slide, "8 thanh ghi CSR cốt lõi:", Inches(0.3), Inches(1.0), Inches(4.4), Inches(0.35),
         font_size=Pt(16), bold=True, color=BK_BLUE)
csrs = [
    ("mstatus", "MIE/MPIE: bật/tắt ngắt toàn cục"),
    ("mtvec",   "Địa chỉ vector ISR"),
    ("mepc",    "PC lệnh bị gián đoạn"),
    ("mcause",  "Mã nguyên nhân trap (bit31=Interrupt)"),
    ("mie",     "Enable từng nguồn ngắt (MSIE/MTIE/MEIE)"),
    ("mip",     "Pending status các nguồn ngắt"),
    ("mtval",   "Giá trị phụ khi trap (bad address)"),
    ("mscratch","Scratch register cho ISR"),
]
for idx, (name, desc) in enumerate(csrs):
    y = Inches(1.4 + idx * 0.65)
    add_rect(slide, Inches(0.3), y, Inches(1.3), Inches(0.55),
             fill_rgb=BK_BLUE)
    add_text(slide, name, Inches(0.3), y + Inches(0.05), Inches(1.3), Inches(0.45),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, Inches(1.65), y + Inches(0.05), Inches(2.9), Inches(0.45),
             font_size=Pt(13), color=DARK_GRAY)

# Trap flow
add_text(slide, "Quy trình Trap (6 bước phần cứng):", Inches(5.0), Inches(1.0), Inches(4.8), Inches(0.35),
         font_size=Pt(16), bold=True, color=BK_BLUE)
trap_steps = [
    "① Nhận diện ngắt hợp lệ: irq & MIE & MEIE",
    "② Flush IF/ID (precise exception model)",
    "③ Lưu PC → mepc,  mcause ← cause code",
    "④ mstatus: MPIE←MIE, MIE←0 (cấm ngắt lồng)",
    "⑤ PC ← mtvec  (nhảy đến ISR)",
    "⑥ Lệnh MRET: MIE←MPIE, PC←mepc",
]
for idx, step in enumerate(trap_steps):
    y = Inches(1.45 + idx * 0.9)
    add_rect(slide, Inches(5.0), y, Inches(4.7), Inches(0.78),
             fill_rgb=LIGHT_BG if idx % 2 == 0 else WHITE,
             line_rgb=BK_BLUE, line_w=Pt(0.5))
    add_text(slide, step, Inches(5.1), y + Inches(0.05), Inches(4.5), Inches(0.68),
             font_size=Pt(14), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — UVM Theory
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Phương pháp luận UVM (Universal Verification Methodology)", "3")
slide_footer(slide, 8)

add_text(slide, "UVM IEEE 1800.2 — Tiêu chuẩn công nghiệp cho Functional Verification",
         Inches(0.3), Inches(0.95), Inches(9.4), Inches(0.35),
         font_size=Pt(15), italic=True, color=DARK_GRAY)

# UVM Architecture diagram (text-based)
components = [
    ("uvm_test", Inches(4.2), Inches(1.4), BK_BLUE),
    ("uvm_env", Inches(4.2), Inches(2.2), RGBColor(0x00,0x80,0xC0)),
    ("uvm_agent", Inches(1.3), Inches(3.1), ORANGE),
    ("uvm_scoreboard", Inches(5.0), Inches(3.1), RGBColor(0x80,0x00,0x80)),
    ("uvm_sequencer", Inches(0.3), Inches(4.2), RGBColor(0x00,0x96,0x3A)),
    ("uvm_driver", Inches(1.7), Inches(4.2), ORANGE),
    ("uvm_monitor", Inches(3.1), Inches(4.2), RGBColor(0x00,0x70,0xB0)),
    ("Coverage", Inches(7.3), Inches(4.2), RGBColor(0x60,0x00,0x90)),
]
for label, x, y, color in components:
    add_rect(slide, x, y, Inches(1.7), Inches(0.6), fill_rgb=color)
    add_text(slide, label, x, y + Inches(0.08), Inches(1.7), Inches(0.44),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Description
desc_items = [
    ("Constrained Random Verification (CRV): sinh lệnh Assembly ngẫu nhiên hợp lệ", 0),
    ("Scoreboard với Reference Model SW: tự động so sánh kết quả RTL vs mô hình", 0),
    ("Functional Coverage: đo mức độ kích hoạt tất cả nhánh quan trọng của thiết kế", 0),
    ("Self-checking: không cần engineer kiểm tra thủ công từng kết quả", 0),
]
add_bullet_box(slide, desc_items,
               Inches(0.3), Inches(5.35), Inches(9.4), Inches(2.0),
               font_size=Pt(15))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Hardware Architecture Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Thiết kế phần cứng: Kiến trúc tổng quan SoC", "4")
slide_footer(slide, 9)

add_text(slide, "Top-level: soc_top.sv  →  riscv_datapath  +  riscv_control  +  csr  +  IMEM/DMEM",
         Inches(0.3), Inches(0.92), Inches(9.4), Inches(0.35),
         font_size=Pt(14), italic=True, color=DARK_GRAY)

add_image(slide, "3_1.png", Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.9))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Datapath
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Thiết kế phần cứng: Datapath 5 tầng", "4")
slide_footer(slide, 10)

add_image(slide, "3.2.png", Inches(0.1), Inches(0.95), Inches(9.8), Inches(6.3))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EX Stage: ALU + M-Unit
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "EX Stage: Khối ALU và M-Unit (Nhân/Chia)", "4")
slide_footer(slide, 11)

# Left: ALU description
add_text(slide, "Khối ALU", Inches(0.3), Inches(1.0), Inches(3.8), Inches(0.38),
         font_size=Pt(18), bold=True, color=BK_BLUE)
alu_items = [
    ("11 phép toán: ADD/SUB/SLL/SLT/SLTU/XOR/SRL/SRA/OR/AND/PASS", 0),
    ("Forwarding MUX: chọn dữ liệu từ MEM/WB/RF", 0),
    ("Input: (RS1 hoặc PC) × (RS2 hoặc Immediate)", 0),
]
add_bullet_box(slide, alu_items, Inches(0.3), Inches(1.45), Inches(3.8), Inches(1.8),
               font_size=Pt(15))

# M-Unit
add_text(slide, "Khối M-Unit (RV32M Extension)", Inches(0.3), Inches(3.35), Inches(4.0), Inches(0.38),
         font_size=Pt(18), bold=True, color=ORANGE)
munit_items = [
    ("FSM 5 trạng thái: IDLE→PREPARE→(DIV_LOOP×32)→FIX_SIGN→DONE", 0),
    ("MUL: 2 chu kỳ  |  DIV/REM: 35 chu kỳ (Restoring Division)", 0),
    ("Interlock: m_unit_busy_o → Hazard Unit → stall IF/ID/EX", 0),
    ("Chia cho 0: DIV→0xFFFFFFFF, DIVS(MIN_INT,−1)→0x80000000", 0),
]
add_bullet_box(slide, munit_items, Inches(0.3), Inches(3.8), Inches(4.0), Inches(2.6),
               font_size=Pt(14), color=DARK_GRAY)

# Right: FSM diagram
add_image(slide, "3.3.drawio.png", Inches(4.5), Inches(1.0), Inches(5.3), Inches(6.2))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Forwarding Unit
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Data Forwarding — Giải quyết RAW Hazard không dùng Stall", "4")
slide_footer(slide, 12)

fwd_items = [
    ("Ưu tiên MEM→EX (sel=2'b01): kết quả ALU tầng MEM — 'mới nhất'", 0),
    ("WB→EX (sel=2'b10): kết quả ghi RF tầng WB — 'cũ hơn'", 0),
    ("Điều kiện: RegWrite=1, rd≠x0, rd khớp rs1/rs2 tầng EX", 0),
    ("Forwarding cho cả rs1 AND rs2 (bao gồm store_data của STORE)", 0),
]
add_bullet_box(slide, fwd_items, Inches(0.3), Inches(1.0), Inches(4.2), Inches(2.5),
               font_size=Pt(16))

# Example code
add_rect(slide, Inches(0.3), Inches(3.6), Inches(4.0), Inches(2.7),
         fill_rgb=RGBColor(0x1E,0x1E,0x1E))
code = (
    "// EX→EX forwarding (sel=2'b01)\n"
    "ADDI x1, x0, 10   // x1=10 @ MEM\n"
    "ADD  x2, x1, x1   // ← forward x1\n\n"
    "// MEM→EX forwarding (sel=2'b10)\n"
    "ADDI x1, x0, 10   // x1=10 @ WB\n"
    "NOP                // 1 buffer\n"
    "ADD  x2, x1, x3   // ← forward x1"
)
add_text(slide, code, Inches(0.4), Inches(3.7), Inches(3.8), Inches(2.5),
         font_size=Pt(13), color=RGBColor(0xD4,0xD4,0xD4))

# Forwarding diagram
add_image(slide, "3.4.drawio.png", Inches(4.5), Inches(0.95), Inches(5.3), Inches(6.3))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — UVM Environment
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Môi trường kiểm thử UVM — Kiến trúc tổng thể", "5")
slide_footer(slide, 13)

add_image(slide, "4.1.drawio.png", Inches(0.1), Inches(0.92), Inches(6.5), Inches(5.5))

uvm_items = [
    ("Sequence Item (Transaction)", 0),
    ("Lệnh 32-bit + ràng buộc opcode/register", 1),
    ("Driver — IMEM/DMEM Responder", 0),
    ("Đáp ứng fetch IF; cấp data cho DMEM", 1),
    ("Monitor — giám sát RF write-back", 0),
    ("Lấy mẫu rf_we + wdata tại WB stage", 1),
    ("Scoreboard + Reference Model", 0),
    ("So sánh RTL vs SW model tự động", 1),
    ("Coverage Model — 199 bins", 0),
    ("opcode, hazard path, fwd_sel, trap", 1),
]
add_bullet_box(slide, uvm_items,
               Inches(6.7), Inches(1.0), Inches(3.1), Inches(6.2),
               font_size=Pt(14))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Results: Unit Tests
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Kết quả kiểm thử mức đơn vị (Unit Test)", "6")
slide_footer(slide, 14)

units = [
    ("ALU", "1000 test vectors PASS\n11 phép toán, biên: overflow/underflow, SRA(MIN_INT), SLTU corner", GREEN_OK),
    ("M-Unit", "500 mẫu ngẫu nhiên PASS\nMUL=2cy, DIV=35cy ✓  |  Chia/0→0xFFFF ✓  |  MIN_INT÷(−1)→0x8000 ✓", GREEN_OK),
    ("Decoder", "47 RV32IM + 9 system instr PASS\nĐủ 6 dạng Immediate; illegal_instr=1 đúng khi opcode không hợp lệ", GREEN_OK),
    ("LSU", "LB/LH/LW/LBU/LHU/SB/SH/SW PASS\nMisaligned LH/LW/SH → lsu_err=1 ✓  |  Handshake 1–5cy delay ✓", GREEN_OK),
]

for idx, (name, result, color) in enumerate(units):
    y = Inches(1.05 + idx * 1.55)
    add_rect(slide, Inches(0.3), y, Inches(1.4), Inches(1.35), fill_rgb=color)
    add_text(slide, name, Inches(0.3), y + Inches(0.35), Inches(1.4), Inches(0.65),
             font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.75), y, Inches(8.0), Inches(1.35),
             fill_rgb=LIGHT_BG, line_rgb=color, line_w=Pt(1))
    add_text(slide, result, Inches(1.9), y + Inches(0.1), Inches(7.7), Inches(1.15),
             font_size=Pt(14), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — System Test: Forwarding + Load-Use
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Kết quả hệ thống: Forwarding & Load-Use Stall", "6")
slide_footer(slide, 15)

# Forwarding result
add_rect(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(2.6),
         fill_rgb=RGBColor(0xE8,0xFF,0xED), line_rgb=GREEN_OK, line_w=Pt(1))
add_text(slide, "✓ Data Forwarding", Inches(0.4), Inches(1.05), Inches(4.0), Inches(0.45),
         font_size=Pt(18), bold=True, color=GREEN_OK)
add_text(slide,
         "• EX→EX (MEM forwarding, sel=2'b01): x2=20 đúng, KHÔNG có Stall\n"
         "• MEM→EX (WB forwarding, sel=2'b10): x2=15 đúng, KHÔNG có Stall\n"
         "• Store forwarding (rs2→store_data) và Branch forwarding: PASS\n"
         "• Waveform chu kỳ 4 & 6: fwd_rs1_sel đúng 2'b01 → 2'b10",
         Inches(0.4), Inches(1.55), Inches(9.0), Inches(1.0),
         font_size=Pt(14), color=DARK_GRAY)

# Load-Use result
add_rect(slide, Inches(0.3), Inches(3.8), Inches(9.4), Inches(2.6),
         fill_rgb=RGBColor(0xFF,0xF3,0xE0), line_rgb=ORANGE, line_w=Pt(1))
add_text(slide, "✓ Load-Use Hazard (Pipeline Stall)", Inches(0.4), Inches(3.85), Inches(5.0), Inches(0.45),
         font_size=Pt(18), bold=True, color=ORANGE)
add_text(slide,
         "• Hazard Detection Unit phát hiện đúng is_load_use = 1\n"
         "• 1 chu kỳ NOP chèn vào giữa LW và ADD  —  x3 nhận giá trị đúng\n"
         "• Negative test: không có phụ thuộc → ctrl_force_stall_id=0 (không false positive)\n"
         "• CPI = 1.08 trong directed test (~50% LW+USE liên tiếp)",
         Inches(0.4), Inches(4.35), Inches(9.0), Inches(1.0),
         font_size=Pt(14), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — System Test: Branch + Interrupt
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Kết quả hệ thống: Branch Flush & Ngắt phần cứng", "6")
slide_footer(slide, 16)

# Branch
add_rect(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(2.6),
         fill_rgb=RGBColor(0xE8,0xFF,0xED), line_rgb=GREEN_OK, line_w=Pt(1))
add_text(slide, "✓ Control Hazard — Branch/Jump Flush", Inches(0.4), Inches(1.05), Inches(5.5), Inches(0.45),
         font_size=Pt(18), bold=True, color=GREEN_OK)
add_text(slide,
         "• BEQ branch_taken=1 → ctrl_flush_if_id=1 & ctrl_flush_id_ex=1 ngay chu kỳ đó\n"
         "• Lệnh sai luồng bị xóa: IF/ID và ID/EX → NOP, PC ← địa chỉ đích\n"
         "• JAL lưu PC+4 vào x1 đúng; lệnh bị FLUSH không thực thi",
         Inches(0.4), Inches(1.55), Inches(9.0), Inches(1.0),
         font_size=Pt(14), color=DARK_GRAY)

# Interrupt
add_rect(slide, Inches(0.3), Inches(3.8), Inches(9.4), Inches(2.6),
         fill_rgb=RGBColor(0xFF,0xF3,0xE0), line_rgb=ORANGE, line_w=Pt(1))
add_text(slide, "✓ External Interrupt (M-mode)", Inches(0.4), Inches(3.85), Inches(5.0), Inches(0.45),
         font_size=Pt(18), bold=True, color=ORANGE)
add_text(slide,
         "• irq_ext_i=1 (sau khi MEIE & MIE đã bật) → ctrl_trap_valid=1, mcause=11\n"
         "• PC ← mtvec(0x40), mepc lưu đúng địa chỉ lệnh bị gián đoạn\n"
         "• ISR thực thi: x5=0xBB ✓  |  MRET: MIE←MPIE, PC←mepc — CPU tiếp tục luồng chính\n"
         "• Độ trễ ngắt: 3–5 chu kỳ kể từ khi irq được xác nhận",
         Inches(0.4), Inches(4.35), Inches(9.0), Inches(1.0),
         font_size=Pt(14), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Stress Test + Coverage
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Kết quả tổng hợp: Stress Test & Functional Coverage", "6")
slide_footer(slide, 17)

# Summary numbers
metrics = [
    ("10,000+", "Lệnh ngẫu nhiên\nkiểm thử", BK_BLUE),
    ("100%", "Scoreboard\nPASS", GREEN_OK),
    ("100%", "Functional\nCoverage", GREEN_OK),
    ("199", "Coverage\nbins đạt", RGBColor(0x00,0x70,0xB0)),
]
for idx, (num, label, color) in enumerate(metrics):
    x = Inches(0.5 + idx * 2.3)
    add_rect(slide, x, Inches(1.05), Inches(2.0), Inches(1.5), fill_rgb=color)
    add_text(slide, num, x, Inches(1.1), Inches(2.0), Inches(0.85),
             font_size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, Inches(1.95), Inches(2.0), Inches(0.6),
             font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

# Stress test details
add_text(slide, "Cấu hình Stress Test:", Inches(0.3), Inches(2.75), Inches(4.5), Inches(0.38),
         font_size=Pt(16), bold=True, color=BK_BLUE)
stress_items = [
    ("50% lệnh ALU — tạo phụ thuộc RAW liên tiếp", 0),
    ("17% LOAD-USE — stress Load-Use Stall logic", 0),
    ("17% BRANCH — stress Branch Flush", 0),
    ("16% STORE — test data memory + forwarding", 0),
    ("Thanh ghi x1–x5 (tối đa trùng địa chỉ → forwarding cao)", 0),
    ("10 lần chạy độc lập × 100 lệnh/lần = 1,000 lệnh/run × 10 seeds", 0),
]
add_bullet_box(slide, stress_items, Inches(0.3), Inches(3.2), Inches(4.5), Inches(3.8),
               font_size=Pt(14))

# Coverage details
add_text(slide, "Coverage Model:", Inches(5.2), Inches(2.75), Inches(4.5), Inches(0.38),
         font_size=Pt(16), bold=True, color=BK_BLUE)
cov_items = [
    ("Tất cả 47 lệnh RV32IM + 9 system instr", 0),
    ("Forwarding paths: MEM/WB → EX (rs1 & rs2)", 0),
    ("Branch taken / not taken", 0),
    ("Trap: Software / Timer / External interrupt", 0),
    ("Pipeline stall (Load-Use, M-Unit busy)", 0),
    ("Kết quả: 100% — 199/199 bins hit", 0),
]
add_bullet_box(slide, cov_items, Inches(5.2), Inches(3.2), Inches(4.5), Inches(3.8),
               font_size=Pt(14))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Bugs found during debug
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Bugs phát hiện & sửa trong quá trình Debug", "6")
slide_footer(slide, 18)

add_text(slide, "UVM Scoreboard đã phát hiện và giúp sửa nhiều lỗi RTL quan trọng:",
         Inches(0.3), Inches(0.95), Inches(9.4), Inches(0.38),
         font_size=Pt(15), italic=True, color=DARK_GRAY)

bugs = [
    ("Priority Forwarding sai",
     "MUX chọn WB thay vì MEM khi cả hai cùng match → x2 = giá trị cũ\nFix: Ưu tiên MEM (sel=2'b01) trước WB (sel=2'b10) trong Forwarding Unit",
     RED_ERR),
    ("CSR trap: MIE xóa muộn 1 cycle",
     "MIE không bị xóa ngay khi trap → interrupt lồng nếu 2 ngắt cùng lúc\nFix: mstatus_mie ← 0 trong cùng cycle với trap_valid_i = 1",
     RED_ERR),
    ("Load-Use Stall thiếu điều kiện",
     "Stall kể cả khi rd = x0 (x0 luôn =0, không cần stall)\nFix: Thêm điều kiện hz_ex_rd != 0 vào is_load_use logic",
     RED_ERR),
    ("M-Unit: DIV chưa xử lý MIN_INT÷(−1)",
     "Kết quả sai do 0x80000000 ÷ 0xFFFFFFFF tràn số nguyên 32-bit\nFix: Phát hiện corner case và trả về 0x80000000 trong FIX_SIGN state",
     RED_ERR),
]

for idx, (bug, fix, color) in enumerate(bugs):
    y = Inches(1.45 + idx * 1.45)
    add_rect(slide, Inches(0.3), y, Inches(9.4), Inches(1.3),
             fill_rgb=RGBColor(0xFF,0xF0,0xF0), line_rgb=RED_ERR, line_w=Pt(0.75))
    add_text(slide, f"🐛  {bug}", Inches(0.4), y + Inches(0.05), Inches(9.0), Inches(0.38),
             font_size=Pt(14), bold=True, color=RED_ERR)
    add_text(slide, fix, Inches(0.5), y + Inches(0.45), Inches(8.8), Inches(0.8),
             font_size=Pt(13), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Kết luận", "7")
slide_footer(slide, 19)

conclude_items = [
    ("Thiết kế hoàn chỉnh lõi RV32IM pipeline 5 tầng bằng SystemVerilog", 0),
    ("47 lệnh RV32IM + 9 lệnh hệ thống CSR; tuân thủ RISC-V spec", 1),
    ("Forwarding Unit, Hazard Detection Unit, M-Unit FSM (MUL=2cy, DIV=35cy)", 1),
    ("CSR + Interrupt M-mode đầy đủ (8 CSR, 3 nguồn ngắt, trap/MRET)", 1),
    ("Xây dựng môi trường UVM chuẩn công nghiệp", 0),
    ("Constrained Random Verification · Scoreboard + RefModel · Coverage 199 bins", 1),
    ("10,000+ lệnh ngẫu nhiên — Scoreboard 100% PASS", 1),
    ("Functional Coverage 100% (tất cả 199 bins được kích hoạt)", 1),
    ("Điểm khác biệt so với các đề tài tương tự trong nước", 0),
    ("Giao thức Valid/Ready chuẩn AXI-like (tiền đề cho AXI-4 bus)", 1),
    ("UVM Constrained Random thay vì Directed Testbench truyền thống", 1),
    ("Hạn chế: Chưa FPGA synthesis  ·  Chưa PLIC  ·  Coverage cross chưa đầy đủ", 0),
]
add_bullet_box(slide, conclude_items, Inches(0.4), Inches(1.0), Inches(9.2), Inches(6.2),
               font_size=Pt(17))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Future Work
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
slide_header(slide, "Hướng phát triển", "7")
slide_footer(slide, 20)

future_items = [
    ("Phần cứng — Hiệu năng & Tài nguyên", 0),
    ("Tổng hợp FPGA (DE10-Standard): đo Fmax, LUT/FF/BRAM thực tế", 1),
    ("Branch Predictor (BTB + BHT 2-bit): giảm penalty từ 2 xuống 0 chu kỳ (90%)", 1),
    ("Tích hợp AXI-4 Lite Bus: kết nối IP vào SoC chuẩn công nghiệp", 1),
    ("Phần cứng — Mở rộng tính năng", 0),
    ("Physical Memory Protection (PMP): phân quyền bộ nhớ theo chuẩn RISC-V", 1),
    ("S-mode + Sv32 MMU: hỗ trợ hệ điều hành nhúng (FreeRTOS/Linux)", 1),
    ("C-Extension (16-bit instructions): tăng mật độ code", 1),
    ("Verification — Nâng cao độ tin cậy", 0),
    ("Formal Verification (SVA): chứng minh x0=0 bất biến, FSM no-deadlock", 1),
    ("Cross-coverage: ngắt trong khi Load-Use stall, ngắt liên tiếp", 1),
    ("RISC-V formal test suite (riscv-tests) để so sánh với chuẩn tham chiếu", 1),
]
add_bullet_box(slide, future_items, Inches(0.4), Inches(1.0), Inches(9.2), Inches(6.2),
               font_size=Pt(16))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, W, H, fill_rgb=BK_BLUE)
add_text(slide, "XIN CẢM ƠN QUÝ THẦY CÔ VÀ HỘI ĐỒNG!",
         Inches(0), Inches(2.0), W, Inches(1.2),
         font_size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Rất mong nhận được câu hỏi và góp ý từ Hội đồng.",
         Inches(0), Inches(3.4), W, Inches(0.6),
         font_size=Pt(20), color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Phạm Tiến Khoa  —  MSSV: 2010344\n"
         "Khoa Điện – Điện tử, Trường Đại học Bách Khoa TP.HCM",
         Inches(0), Inches(4.5), W, Inches(0.9),
         font_size=Pt(16), color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
slide_footer(slide, 21)

# ──────────────────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"✓ Saved: {OUT_FILE}  ({len(prs.slides)} slides)")
